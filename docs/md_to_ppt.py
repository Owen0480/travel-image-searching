# -*- coding: utf-8 -*-
"""
통합테스트_시나리오.md → PowerPoint(.pptx) 자동 생성
실행: python md_to_ppt.py  (또는 python docs/md_to_ppt.py)
필요: pip install python-pptx
"""
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx가 필요합니다. 터미널에서 실행: pip install python-pptx")
    exit(1)

MD_PATH = Path(__file__).resolve().parent / "통합테스트_시나리오.md"
OUT_PATH = Path(__file__).resolve().parent / "통합테스트_시나리오.pptx"


def parse_md(path):
    text = path.read_text(encoding="utf-8")
    sections = []
    current = None
    for line in text.split("\n"):
        if line.startswith("## "):
            if current:
                sections.append(current)
            title = line[3:].strip()
            current = {"title": title, "body": [], "table": None}
        elif current is not None:
            if line.strip() == "---":
                continue
            if line.startswith("|"):
                if current.get("table") is None:
                    current["table"] = []
                row = [c.strip() for c in line.split("|")[1:-1]]
                if row and not all(c.replace("-", "").strip() == "" for c in row):
                    current["table"].append(row)
            else:
                stripped = line.strip()
                if stripped and not stripped.startswith("|"):
                    current["body"].append(stripped)
    if current:
        sections.append(current)
    return sections


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(0.8))
        tx2.text_frame.paragraphs[0].text = subtitle
        tx2.text_frame.paragraphs[0].font.size = Pt(18)
        tx2.text_frame.paragraphs[0].font.color.rgb = RgbColor(0x66, 0x66, 0x66)


def add_content_slide(prs, section):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = section["title"]
    p.font.size = Pt(24)
    p.font.bold = True

    top = 1.1
    if section.get("body"):
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(2))
        tf = body_box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(section["body"][:8]):
            if line.startswith("- **"):
                m = re.match(r"-\s*\*\*(.+?)\*\*:\s*(.+)", line)
                if m:
                    line = f"• {m.group(1)}: {m.group(2)}"
                else:
                    line = "• " + re.sub(r"\*\*(.+?)\*\*", r"\1", line)
            elif line.startswith("- "):
                line = "• " + line[2:]
            elif line.startswith("| "):
                continue
            p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
            p.text = line[:120] + ("..." if len(line) > 120 else "")
            p.font.size = Pt(12)
            p.space_after = Pt(4)
        top = 2.2

    if section.get("table") and len(section["table"]) > 0:
        rows, cols = len(section["table"]), len(section["table"][0])
        col_w = 9.0 / cols
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(top), Inches(9), Inches(0.35 * rows)).table
        for r, row in enumerate(section["table"]):
            for c, cell_text in enumerate(row):
                if c < cols:
                    cell = table.cell(r, c)
                    cell.text = cell_text[:50] + (".." if len(cell_text) > 50 else "")
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(10)
                    if r == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RgbColor(0x44, 0x72, 0xC4)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.color.rgb = RgbColor(0xFF, 0xFF, 0xFF)
                            paragraph.font.bold = True


def main():
    sections = parse_md(MD_PATH)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 표지
    add_title_slide(prs, "통합 테스트 시나리오", "여행 숙박 추천 AI 모델")

    # 추천 구성 섹션은 마지막에
    normal = [s for s in sections if "PPT에 넣을 때" not in s["title"] and "추천 슬라이드" not in s["title"]]
    recommend = [s for s in sections if "PPT에 넣을 때" in s["title"] or "추천 슬라이드" in s["title"]]

    for section in normal:
        add_content_slide(prs, section)

    for section in recommend:
        add_content_slide(prs, section)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"저장 완료: {OUT_PATH}")


if __name__ == "__main__":
    main()
